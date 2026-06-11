#!/usr/bin/env python3
"""Replace the embedding-projection figure on the existing PCA/t-SNE slide in
dplm.pptx in place (slide #314) with the regenerated, type-grouped-colour PNG,
and refresh the annotation's colour sentence.

Finds the slide by title (robust to index shifts), removes its single picture,
re-adds the new PNG centred in the same slot (slot recomputed from the new image
aspect ratio), and rewrites the colour paragraph of the annotation box. Run only
with the deck closed (a `~$dplm.pptx` lock means PowerPoint has it open).
"""
from __future__ import annotations
import datetime
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

PPTX = "/Users/soldatmat/Documents/terpene_synthases/presentation/dplm.pptx"
FIG = ("/Users/soldatmat/Documents/terpene_synthases/dplm/run/class_predictor/"
       "slide306_eval/train_tps_pca_tsne.png")
TITLE_PREFIX = "Training TPSs in DPLM-150m embedding space"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)  # same slot as the original append

ANNOTATION = (
    "Same DPLM-150m mean\n"
    "embeddings the kNN first-\n"
    "cyclization classifier uses\n"
    "(z-scored per dim).\n\n"
    "Colour = first-cyclization\n"
    "class; HUE FAMILY =\n"
    "substrate type (blue sesqui,\n"
    "green mono, red di, purple\n"
    "sester, brown sterol). Legend\n"
    "grouped by type.\n\n"
    "Classes do NOT form 22\n"
    "globally separated blobs;\n"
    "instead they cluster\n"
    "LOCALLY (e.g. class 11\n"
    "mono splits off cleanly,\n"
    "t-SNE shows many tight\n"
    "same-class pockets) with\n"
    "broad overlap between\n"
    "related classes. That local\n"
    "structure is exactly what a\n"
    "k=3-9 neighbour vote relies\n"
    "on — global separation is\n"
    "neither present nor needed."
)


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)

    target, tidx = None, None
    for i, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            if sh.has_text_frame and sh.name == "TextovéPole 11" \
                    and sh.text_frame.text.startswith(TITLE_PREFIX):
                target, tidx = sl, i
                break
        if target:
            break
    if target is None:
        raise SystemExit(f"could not find the embedding slide (title prefix {TITLE_PREFIX!r})")

    pics = [s for s in target.shapes if s.shape_type == 13]
    if len(pics) != 1:
        raise SystemExit(f"expected exactly 1 picture on the slide, found {len(pics)}")
    pics[0].element.getparent().remove(pics[0].element)

    l, t, w, h = fit_into_slot(FIG, PIC_SLOT)
    target.shapes.add_picture(FIG, Inches(l), Inches(t), Inches(w), Inches(h))

    for shp in target.shapes:
        if shp.has_text_frame and shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, ANNOTATION)

    prs.save(PPTX)
    print(f"replaced figure on slide {tidx + 1}; slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
