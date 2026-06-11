#!/usr/bin/env python3
"""Single serialized deck update for dplm.pptx (run with PowerPoint CLOSED):

  1. Replace the figure on the existing embedding slide (#314) with the
     regenerated train_tps_pca_tsne.png (legend now ordered by carbon number).
  2. Append 5 generated-sequence overlay slides (class-0 overview + per-arch
     CA/PRE/MINI, and class-1 mini), modelled on the embedding slide's layout
     (template slide index 232: title / slide-number / annotation + big picture).

One Presentation open + one save, so there is no risk of concurrent writers.
"""
from __future__ import annotations
import datetime
import shutil
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = "/Users/soldatmat/Documents/terpene_synthases"
PPTX = f"{ROOT}/presentation/dplm.pptx"
FIGDIR = f"{ROOT}/dplm/run/class_predictor/slide306_eval"

TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

EMBED_TITLE_PREFIX = "Training TPSs in DPLM-150m embedding space"
TRAIN_FIG = f"{FIGDIR}/train_tps_pca_tsne.png"

# (figure, title, annotation) for the 5 appended overlay slides
OVERLAYS = [
    (f"{FIGDIR}/gen_overlay_class0_overview.png",
     "Generated seqs (class 0) in train TPS embedding space – baseline vs best-per-architecture",
     "Generated seqs\nconditioned on class 0\n(target = gold cloud).\n\n"
     "baseline = black X\n(unconditional).\n\n"
     "Best per architecture\n(by kNN class-0\nfidelity): CA FT,rand /\nPRE QVKO / MINI QVK.\n\n"
     "All land in the dense\ncentral TPS region\noverlapping class 0;\n"
     "conditional models do\nNOT separate from\nbaseline in 2D —\nconsistent with the\n"
     "weak measured fidelity\n(frac class0: baseline\n0.46 > every\nconditional run)."),

    (f"{FIGDIR}/gen_overlay_class0_CA.png",
     "Generated seqs (class 0) – cross-attention (CA) variants vs baseline",
     "Class-0 generated,\ncross-attention (CA)\nvariants vs baseline.\n\n"
     "target class 0 = gold;\nbaseline = black X.\n\n"
     "FT,rand (0.26) and\nFT,orig (0.22) sit\nnearest the class-0\ncloud; ALLadap orig\n"
     "(0.04) collapses with\nbaseline into the\ngeneric TPS core.\n\n"
     "(frac = kNN class-0\nfidelity.)"),

    (f"{FIGDIR}/gen_overlay_class0_PRE.png",
     "Generated seqs (class 0) – prepend variants vs baseline",
     "Class-0 generated,\nprepend variants vs\nbaseline.\n\n"
     "target class 0 = gold;\nbaseline = black X.\n\n"
     "QVKO / V29 / V all\ncluster tightly in the\ncentral TPS region,\nlargely on top of\n"
     "baseline — the prepend\nfamily shows the\nweakest class-0\nfidelity (<=0.10)."),

    (f"{FIGDIR}/gen_overlay_class0_MINI.png",
     "Generated seqs (class 0) – mini-CA variants vs baseline",
     "Class-0 generated,\nmini-CA variants vs\nbaseline.\n\n"
     "target class 0 = gold;\nbaseline = black X.\n\n"
     "QVK (0.30) is the best\nconditional overall and\nsits nearest the\nclass-0 cloud; V15-28\n"
     "(0.22) follows; ltm0\n(0.04, no LoRA)\ncollapses with\nbaseline."),

    (f"{FIGDIR}/gen_overlay_class1_MINI.png",
     "Generated seqs (class 1) – mini-CA variants (CA/prepend unavailable: pre-merge)",
     "Generated seqs\nconditioned on class 1\n(target = gold cloud).\n\n"
     "baseline = black X\n(unconditional; class 1\nis a target it never\nproduces).\n\n"
     "MINI variants only —\nCA / prepend run_1 are\npre-merge 24-class and\nfail to load.\n\n"
     "None land on the\nclass-1 cloud (fidelity\n~0.00); ltm0 forms its\nown off-target lobe.\n"
     "Matches the\n'conditioning too weak\nto steer' verdict."),
]


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def add_overlay_slide(prs, template, fig, title, annotation):
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:  # skip the template's picture
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(fig, PIC_SLOT)
    target.shapes.add_picture(fig, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, title)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, annotation)


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    for f in [TRAIN_FIG] + [o[0] for o in OVERLAYS]:
        if not Path(f).exists():
            raise SystemExit(f"missing figure: {f}")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)

    # ---- 1. swap the image on the existing embedding slide (#314) ----
    embed = None
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and sh.name == TITLE_SHAPE_NAME \
                    and sh.text_frame.text.startswith(EMBED_TITLE_PREFIX):
                embed = sl
                break
        if embed:
            break
    if embed is None:
        raise SystemExit("could not find the embedding slide to update")
    pics = [s for s in embed.shapes if s.shape_type == 13]
    if len(pics) != 1:
        raise SystemExit(f"expected 1 picture on embedding slide, found {len(pics)}")
    pics[0].element.getparent().remove(pics[0].element)
    l, t, w, h = fit_into_slot(TRAIN_FIG, PIC_SLOT)
    embed.shapes.add_picture(TRAIN_FIG, Inches(l), Inches(t), Inches(w), Inches(h))
    print("swapped embedding-slide figure (reordered legend)")

    # ---- 2. append the 5 overlay slides ----
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    start = len(prs.slides)
    for fig, title, ann in OVERLAYS:
        add_overlay_slide(prs, template, fig, title, ann)
    print(f"appended {len(OVERLAYS)} overlay slides at #{start + 1}..#{len(prs.slides)}")

    prs.save(PPTX)
    print(f"saved. total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
