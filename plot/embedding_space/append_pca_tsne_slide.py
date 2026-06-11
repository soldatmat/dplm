#!/usr/bin/env python3
"""Append the PCA + t-SNE embedding-projection slide to dplm.pptx.

Mirrors slide 306's layout (template slide index 232: title / slide-number /
annotation text boxes + a single big picture). Run ONLY when PowerPoint has the
deck closed (a `~$dplm.pptx` lock file means it is open — appending then risks
clobbering unsaved work / corrupting the file).
"""
from __future__ import annotations
import datetime
import shutil
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

PPTX = "/Users/soldatmat/Documents/terpene_synthases/presentation/dplm.pptx"
FIG = ("/Users/soldatmat/Documents/terpene_synthases/dplm/run/class_predictor/"
       "slide306_eval/train_tps_pca_tsne.png")
TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

TITLE = ("Training TPSs in DPLM-150m embedding space "
         "– PCA & t-SNE, coloured by first-cyclization class")
ANNOTATION = (
    "Same DPLM-150m mean\n"
    "embeddings the kNN first-\n"
    "cyclization classifier uses\n"
    "(z-scored per dim).\n\n"
    "Colour = 22 first-\n"
    "cyclization classes (legend\n"
    "shows id + substrate type).\n\n"
    "Classes do NOT form 22\n"
    "globally separated blobs;\n"
    "instead they cluster\n"
    "LOCALLY (e.g. class 11\n"
    "mono splits off cleanly,\n"
    "t-SNE shows many tight\n"
    "same-class pockets) with\n"
    "broad overlap between\n"
    "related classes. That local\n"
    "structure is exactly what a\n"
    "k=3-9 neighbour vote relies\n"
    "on — global separation is\n"
    "neither present nor needed."
)


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(FIG, PIC_SLOT)
    target.shapes.add_picture(FIG, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, TITLE)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, ANNOTATION)
    prs.save(PPTX)
    print(f"saved. slide count: {len(prs.slides)} (new slide at index {len(prs.slides)-1})")


if __name__ == "__main__":
    main()
