#!/usr/bin/env python3
"""Append the tps_eval visualization-suite slides (new modalities + UMAP
companions) to presentation/dplm.pptx, modelled on the ESM embedding slide
layout (template slide index 232). Figures are produced by tps_eval's
src/visualization/. Run with PowerPoint CLOSED. Add/remove SLIDES entries as
more maps (e.g. EE) become available.
"""
from __future__ import annotations
import datetime
import shutil
import sys
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = "/Users/soldatmat/Documents/terpene_synthases"
PPTX = f"{ROOT}/presentation/dplm.pptx"
VISUALIZATION_DIR = f"{ROOT}/tps_eval/data/visualization"

TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

ALL_SLIDES = {
    "saprot": (f"{VISUALIZATION_DIR}/saprot_map.png",
     "MARTS-DB TPSs — SaProt structure-aware embeddings (PaCMAP / t-SNE / UMAP)",
     "SaProt-650M structure-\naware embeddings (per-\nresidue AA + foldseek-3Di\ntoken), mean-pooled;\nn=984 unique enzymes.\n\n"
     "Colour = first-cyclization\nclass; hue = substrate\ntype.\n\n"
     "Structure-aware PLM\ncleanly separates folds/\nclasses: sterol/OSC (12,\nbrown) isolates, di (reds)\ncluster, sesqui (blues) /\nmono (greens) resolve.\n\n"
     "Strongest new structural\nmodality — folds structure\ninto a continuous embedding\n(fixes the saturated\nfoldseek map)."),

    "esm": (f"{VISUALIZATION_DIR}/esm_umap.png",
     "MARTS-DB TPSs — DPLM-150m ESM embeddings (UMAP)",
     "UMAP companion to the\nexisting PCA / t-SNE map\nof the DPLM-150m mean\nembeddings; n=1349.\n\n"
     "Colour = first-cyclization\nclass (carbon-ordered\nsubstrate-type palette).\n\n"
     "UMAP preserves local +\nsome global structure;\ncompare the cluster layout\nto the PCA/t-SNE version."),

    "seq": (f"{VISUALIZATION_DIR}/seq_umap.png",
     "MARTS-DB TPSs — sequence identity (UMAP, mmseqs)",
     "UMAP on the mmseqs all-\nvs-all sequence-identity\ndistance (1 - fident),\nprecomputed metric;\nn=991 unique enzymes.\n\n"
     "Colour = first-cyclization\nclass.\n\n"
     "UMAP companion to the\nseq-identity PCoA / t-SNE.\nLow TPS identity still\nlimits separation, but UMAP\nrecovers more local cluster\nstructure than the PCoA\nhairball."),

    "struct": (f"{VISUALIZATION_DIR}/struct_umap.png",
     "MARTS-DB TPSs — structural similarity (UMAP, foldseek TM-score)",
     "UMAP on the foldseek all-\nvs-all TM-score distance\n(1 - alntmscore),\nprecomputed; n=1319\nESMFold structures.\n\n"
     "Colour = first-cyclization\nclass.\n\n"
     "UMAP companion to the\nfoldseek PCoA / t-SNE —\nseparates folds better than\nthe saturated-matrix PCoA."),

    "active_site": (f"{VISUALIZATION_DIR}/active_site_umap.png",
     "MARTS-DB TPSs — active-site / cation-residue features (UMAP)",
     "Active-site / cation-\nresidue features → UMAP:\n32-d property + geometry\nprofile of the shell within\n12A of the Mg2+\ncarboxylate cage.\n\n"
     "Class-I fold only (mono /\nsesqui / di / sester); OSC\nclass-12 EXCLUDED (different\nfold, no Mg2+ cage). n=837.\n\n"
     "Colour = first-cyclization\nclass. Modest product-class\nsignal (kNN 0.39 vs 0.27\nbaseline) — partial, not\nclean, separation, as\nexpected cross-fold."),

    "ee_plm": (f"{VISUALIZATION_DIR}/ee_esm1v_map.png",
     "MARTS-DB TPSs — EnzymeExplorer ESM-1v-TPS embeddings (PCA / t-SNE / UMAP)",
     "EnzymeExplorer PLM block:\nESM-1v (650M), TPS-\nfinetuned 'subseq'\ncheckpoint, mean-pooled\n(1280-d); n=991.\n\nColour = first-cyclization\nclass; hue = substrate type.\n\nA second PLM view, distinct\nfrom the DPLM-150m map\n(different model, EE-\nfinetuned). Classes/types\nresolve: sterol (12) + di\n(reds) cluster, sesqui\n(blues) / mono (greens)\nseparate."),

    "ee_domain": (f"{VISUALIZATION_DIR}/ee_domain_map.png",
     "MARTS-DB TPSs — EnzymeExplorer AF-domain-similarity features (PCA / t-SNE / UMAP)",
     "EnzymeExplorer structure/\nfunction block: pairwise\nsimilarities to reference\nfunctional domains (AF/\nESMFold domain segmentation;\nthe domains_subset the\nproduction model uses).\n\nColour = first-cyclization\nclass; hue = substrate type.\n\nThe complementary struct/\nfunction signal not captured\nby the PLM maps — EE's\nactual innovation."),
}

# order to append (override via argv: a subset of the keys above)
ORDER = ["saprot", "esm", "seq", "struct", "active_site"]


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def add_slide(prs, template, fig, title, annotation):
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(fig, PIC_SLOT)
    target.shapes.add_picture(fig, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, title)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, annotation)


def main():
    keys = sys.argv[1:] or ORDER
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    slides = [ALL_SLIDES[k] for k in keys]
    for fig, _, _ in slides:
        if not Path(fig).exists():
            raise SystemExit(f"missing figure: {fig}")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    start = len(prs.slides)
    for fig, title, ann in slides:
        add_slide(prs, template, fig, title, ann)
    prs.save(PPTX)
    print(f"appended {len(slides)} slides ({', '.join(keys)}) at #{start + 1}..#{len(prs.slides)}; total {len(prs.slides)}")


if __name__ == "__main__":
    main()
