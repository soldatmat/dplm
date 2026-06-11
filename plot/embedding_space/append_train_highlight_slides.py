#!/usr/bin/env python3
"""Append the two train-only class-highlight slides (class 0, class 1; no
generated seqs) to presentation/dplm.pptx, modelled on the ESM embedding slide
layout (template slide index 232). Run with PowerPoint CLOSED.
"""
from __future__ import annotations
import datetime
import shutil
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = "/Users/soldatmat/Documents/terpene_synthases"
PPTX = f"{ROOT}/presentation/dplm.pptx"
FIGDIR = f"{ROOT}/dplm/run/class_predictor/slide306_eval"

TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

SLIDES = [
    (f"{FIGDIR}/train_highlight_class0.png",
     "Training TPSs – first-cyclization class 0 highlighted (no generated seqs)",
     "Train-only ESM\nembedding map — same\nlayout & axes as the\nall-class map.\n\n"
     "Gold = class 0\n(target) training TPSs\n(n=227); grey = all\nother training classes\n(n=1122).\n\n"
     "NO generated sequences\n— clean reference for\nwhere the class-0 cloud\nsits, to read alongside\n"
     "the class-0 generated-\noverlay slides.\n\nclass 0 = sesqui (FPP)."),

    (f"{FIGDIR}/train_highlight_class1.png",
     "Training TPSs – first-cyclization class 1 highlighted (no generated seqs)",
     "Train-only ESM\nembedding map — same\nlayout & axes as the\nall-class map.\n\n"
     "Gold = class 1\n(target) training TPSs\n(n=116); grey = all\nother training classes\n(n=1233).\n\n"
     "NO generated sequences\n— clean reference for\nwhere the class-1 cloud\nsits, to read alongside\n"
     "the class-1 generated-\noverlay slide.\n\nclass 1 = sesqui (FPP,\nmarts_M00031)."),
]


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def add_slide(prs, template, fig, title, annotation):
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(fig, PIC_SLOT)
    target.shapes.add_picture(fig, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, title)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, annotation)


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    for fig, _, _ in SLIDES:
        if not Path(fig).exists():
            raise SystemExit(f"missing figure: {fig}")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    start = len(prs.slides)
    for fig, title, ann in SLIDES:
        add_slide(prs, template, fig, title, ann)
    prs.save(PPTX)
    print(f"appended {len(SLIDES)} slides at #{start + 1}..#{len(prs.slides)}; total {len(prs.slides)}")


if __name__ == "__main__":
    main()
