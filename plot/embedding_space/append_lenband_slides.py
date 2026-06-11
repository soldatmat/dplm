#!/usr/bin/env python3
"""Append the three length-band slides to presentation/dplm.pptx:
  - class-0 highlight + 280-420 AA (one-domain) members in blue
  - class-1 highlight + 280-420 AA members in blue
  - all-class map with 280-420 AA points outlined
Modelled on the ESM embedding slide layout (template slide index 232).
Run with PowerPoint CLOSED.
"""
from __future__ import annotations
import datetime
import shutil
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = "/Users/soldatmat/Documents/terpene_synthases"
PPTX = f"{ROOT}/presentation/dplm.pptx"
FIGDIR = f"{ROOT}/dplm/run/class_predictor/slide306_eval"

TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

SLIDES = [
    (f"{FIGDIR}/train_highlight_class0_lenband.png",
     "Training TPSs – class 0 highlighted; 280-420 AA (~1 structural domain) in blue",
     "Train-only ESM map\n(no generated seqs),\nsame axes as the\nall-class map.\n\n"
     "Gold = class 0 (n=227).\nBLUE = the class-0\nsubset with sequence\nlength 280-420 AA\n(~1 structural domain):\n70 of 227.\nGrey = other classes.\n\n"
     "The one-domain (blue)\nmembers form their own\ncluster, well separated\nfrom the longer gold\nclass-mates — single-\ndomain TPSs occupy a\ndistinct region."),

    (f"{FIGDIR}/train_highlight_class1_lenband.png",
     "Training TPSs – class 1 highlighted; 280-420 AA (~1 structural domain) in blue",
     "Train-only ESM map\n(no generated seqs),\nsame axes as the\nall-class map.\n\n"
     "Gold = class 1 (n=116).\nBLUE = the class-1\nsubset with sequence\nlength 280-420 AA\n(~1 structural domain):\n23 of 116.\nGrey = other classes.\n\n"
     "Same pattern as class 0:\nthe one-domain (blue)\nmembers separate from\nthe longer gold members."),

    (f"{FIGDIR}/train_tps_pca_tsne_lenband.png",
     "Training TPSs by first-cyclization class; 280-420 AA (~1 structural domain) outlined",
     "All-class ESM map\n(same layout/axes as\nthe by-class map),\ncolour = first-\ncyclization class\n(hue = substrate type).\n\n"
     "BLACK OUTLINE = the\n321 / 1349 TPSs with\nsequence length\n280-420 AA (~1\nstructural domain).\n\n"
     "Outlined (one-domain)\npoints concentrate in\nthe sesqui (blue) family\nregions; the sterol /\ntriterpene class (12,\nbrown) and much of\nmono are un-outlined\n(longer, multi-domain)."),
]


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def add_slide(prs, template, fig, title, annotation):
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(fig, PIC_SLOT)
    target.shapes.add_picture(fig, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, title)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, annotation)


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    for fig, _, _ in SLIDES:
        if not Path(fig).exists():
            raise SystemExit(f"missing figure: {fig}")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    start = len(prs.slides)
    for fig, title, ann in SLIDES:
        add_slide(prs, template, fig, title, ann)
    prs.save(PPTX)
    print(f"appended {len(SLIDES)} slides at #{start + 1}..#{len(prs.slides)}; total {len(prs.slides)}")


if __name__ == "__main__":
    main()
