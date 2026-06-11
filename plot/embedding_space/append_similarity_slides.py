#!/usr/bin/env python3
"""Append the two similarity-map slides (sequence-identity, structural-TMscore)
to presentation/dplm.pptx, modelled on the ESM embedding slide's layout
(template slide index 232: title / slide-number / annotation + big picture).

Run with PowerPoint CLOSED (a ~$dplm.pptx lock means it's open).
"""
from __future__ import annotations
import datetime
import shutil
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = "/Users/soldatmat/Documents/terpene_synthases"
PPTX = f"{ROOT}/presentation/dplm.pptx"
SIMDIR = f"{ROOT}/dplm/run/class_predictor/slide306_eval/similarity_data"

TEMPLATE_SLIDE_INDEX = 232
LAYOUT_INDEX = 1
TITLE_SHAPE_NAME = "TextovéPole 11"
SLIDE_NUMBER_SHAPE_NAME = "TextovéPole 12"
ANNOTATION_SHAPE_NAME = "TextBox 5"
PIC_SLOT = (0.12, 0.96, 8.48, 6.35)

SLIDES = [
    (f"{SIMDIR}/seq_similarity_map.png",
     "Training TPSs in pairwise SEQUENCE-identity space – PCoA & t-SNE (mmseqs all-vs-all)",
     "2D layout from pairwise\nSEQUENCE identity\n(mmseqs all-vs-all),\nNOT ESM embeddings.\n\n"
     "distance = 1 - fraction\nidentical (fident);\nunreported pairs = max.\n991 unique enzymes.\n\n"
     "Colour = first-cyclization\nclass; hue = substrate\ntype (carbon order).\n\n"
     "Classes overlap heavily —\nTPS sequence identity is\nlow across the family, so\n"
     "sequence space clusters\nby class only weakly\n(t-SNE largely a hairball;\n"
     "sterol/cls 12 brown still\nseparates). PCoA variance\nis low (saturated matrix)."),

    (f"{SIMDIR}/struct_similarity_map.png",
     "Training TPSs in pairwise STRUCTURAL-similarity space – PCoA & t-SNE (foldseek TM-score)",
     "2D layout from pairwise\nSTRUCTURAL similarity\n(foldseek all-vs-all TM-\nscore), NOT ESM embed.\n\n"
     "distance = 1 - foldseek\nalntmscore; unreported\npairs = max. 1319\nESMFold structures.\n\n"
     "Colour = first-cyclization\nclass; hue = substrate\ntype (carbon order).\n\n"
     "Structure separates folds\nMUCH better than seq:\nsterol/cls 12 (brown, OSC\n"
     "fold) isolates cleanly, di\nclasses cluster. Read the\nt-SNE panel — PCoA\n"
     "variance is low (foldseek\nreports only ~100\nneighbours/query)."),
]


def fit_into_slot(img_path, slot):
    sl, st, sw, sh = slot
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio, slot_ratio = iw / ih, sw / sh
    if img_ratio > slot_ratio:
        w, h = sw, sw / img_ratio
    else:
        h, w = sh, sh * img_ratio
    return sl + (sw - w) / 2, st + (sh - h) / 2, w, h


def set_text(text_frame, new_text):
    lines = new_text.split("\n")
    p0 = text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p0.text = lines[0]
    for extra_p in list(text_frame.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def add_slide(prs, template, fig, title, annotation):
    target = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])
    for ph in list(target.placeholders):
        ph.element.getparent().remove(ph.element)
    for shp in template.shapes:
        if shp.shape_type == 13:
            continue
        target.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")
    l, t, w, h = fit_into_slot(fig, PIC_SLOT)
    target.shapes.add_picture(fig, Inches(l), Inches(t), Inches(w), Inches(h))
    for shp in target.shapes:
        if shp.shape_type == 13 or not shp.has_text_frame:
            continue
        if shp.name == TITLE_SHAPE_NAME:
            set_text(shp.text_frame, title)
        elif shp.name == SLIDE_NUMBER_SHAPE_NAME:
            set_text(shp.text_frame, "")
        elif shp.name == ANNOTATION_SHAPE_NAME:
            set_text(shp.text_frame, annotation)


def main():
    if list(Path(PPTX).parent.glob("~$*.pptx")):
        raise SystemExit("ABORT: a ~$*.pptx lock file is present — close PowerPoint first.")
    for fig, _, _ in SLIDES:
        if not Path(fig).exists():
            raise SystemExit(f"missing figure: {fig}")
    backup = f"/tmp/dplm_pptx_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx"
    shutil.copy2(PPTX, backup)
    print(f"backup -> {backup}")
    prs = Presentation(PPTX)
    template = prs.slides[TEMPLATE_SLIDE_INDEX]
    start = len(prs.slides)
    for fig, title, ann in SLIDES:
        add_slide(prs, template, fig, title, ann)
    prs.save(PPTX)
    print(f"appended {len(SLIDES)} slides at #{start + 1}..#{len(prs.slides)}; total {len(prs.slides)}")


if __name__ == "__main__":
    main()
